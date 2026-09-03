"""
generate_ppt.py — Generates the Drift-Sense Phase 2 presentation as .pptx
Run:  pip install python-pptx
      python3 generate_ppt.py
Output: drift_sense_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour Palette ──
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD    = RGBColor(0x1A, 0x22, 0x36)
ACCENT     = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_LT  = RGBColor(0x60, 0xA5, 0xFA)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
RED        = RGBColor(0xEF, 0x44, 0x44)
YELLOW     = RGBColor(0xF5, 0x9E, 0x0B)
WHITE      = RGBColor(0xF1, 0xF5, 0xF9)
GRAY       = RGBColor(0x94, 0xA3, 0xB8)
MUTED      = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT   = RGBColor(0x1E, 0x29, 0x3B)

# ── Helpers ──
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox, tf

def add_para(tf, text, font_size=14, color=WHITE, bold=False, space_before=Pt(4),
             space_after=Pt(2), alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_bullet(tf, text, font_size=13, color=GRAY, bold=False, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.level = level
    p.space_before = Pt(2)
    p.space_after = Pt(2)
    return p

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None,
              header_color=ACCENT, header_text_color=WHITE, cell_color=BG_CARD,
              cell_text_color=GRAY, font_size=10):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]

            cell_data = data[r][c] if r < len(data) and c < len(data[r]) else ""

            # Handle tuples: (text, color, bold)
            if isinstance(cell_data, tuple):
                p.text = str(cell_data[0])
                p.font.color.rgb = cell_data[1] if len(cell_data) > 1 else cell_text_color
                p.font.bold = cell_data[2] if len(cell_data) > 2 else False
            else:
                p.text = str(cell_data)

            p.font.size = Pt(font_size)
            p.font.name = "Calibri"

            if r == 0:  # header
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                p.font.color.rgb = header_text_color
                p.font.bold = True
                p.font.size = Pt(font_size - 1)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_color
                if not isinstance(cell_data, tuple):
                    p.font.color.rgb = cell_text_color

            p.alignment = PP_ALIGN.LEFT
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)

    return table_shape

def slide_header(slide, number, label, title):
    set_slide_bg(slide, BG_DARK)
    # slide number label
    add_textbox(slide, 0.5, 0.25, 4, 0.3, f"{number}  ·  {label}",
                font_size=8, color=MUTED, bold=True)
    # title
    add_textbox(slide, 0.5, 0.5, 9, 0.6, title,
                font_size=26, color=WHITE, bold=True)
    # accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05),
                                   Inches(1.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.shadow.inherit = False

def footer_bar(slide):
    bar = add_rect(slide, 0, 7.1, 10, 0.4, RGBColor(0x0A, 0x0E, 0x1A))
    add_textbox(slide, 0.5, 7.15, 5, 0.25,
                "Drift-Sense  ·  Yash Kothari  ·  BITS Pilani  ·  Semicon India Hackathon",
                font_size=7, color=MUTED)

# ══════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]  # blank

# ═══ SLIDE 1: TITLE ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_DARK)

# Subtle accent glow (rectangle behind title)
glow = add_rect(s, 2.5, 1.5, 5, 3, RGBColor(0x12, 0x1B, 0x30))

# Badge
badge = add_rect(s, 3.2, 1.8, 3.6, 0.35, ACCENT)
add_textbox(s, 3.2, 1.82, 3.6, 0.3, "SEMICON INDIA HACKATHON 2026",
            font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Title
add_textbox(s, 1, 2.5, 8, 0.8, "Drift-Sense",
            font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(s, 1.5, 3.3, 7, 0.7,
            "AI-Powered Navigation-Error Recovery\nfor Wafer Inspection Tools",
            font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# Divider
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.6), Inches(4.2), Inches(0.8), Pt(2))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
line.shadow.inherit = False

# Author
add_textbox(s, 1, 4.5, 8, 0.4, "Yash Kothari",
            font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1, 4.9, 8, 0.4, "BITS Pilani — Electronics & Communication Engineering",
            font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1, 5.3, 8, 0.4, "Applied Materials — Problem Statement 2",
            font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)

footer_bar(s)

# ═══ SLIDE 2: PROBLEM STATEMENT ═══════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "01", "PROBLEM STATEMENT", "Navigation-Error Recovery, Not Template Matching")

# Left column — description
_, tf = add_textbox(s, 0.5, 1.35, 4.5, 2.5, "", font_size=13, color=GRAY)
add_para(tf, "Given a high-resolution reference (1000×1000, 1 nm/px, 100× zoom) "
         "and a wide-search image (1000×1000, ~10 nm/px), locate the reference "
         "pattern's center to sub-pixel precision.", 13, GRAY)
add_para(tf, "", 6, GRAY)
add_para(tf, "Key Reframing", 14, ACCENT_LT, True, Pt(12))
add_para(tf, "The tool knows where it commanded the stage. Drift is a small, "
         "bounded error — not an arbitrary location in the 10 μm field. "
         "A recovery system exploits this prior.", 12, GRAY)

# Right column — cards
card1 = add_rect(s, 5.3, 1.35, 4.2, 1.3, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 5.5, 1.4, 3.8, 1.2, "", font_size=12)
add_para(tf, "Phase 1 Constraints", 12, ACCENT_LT, True)
add_bullet(tf, "Exact 10× zoom ratio (known)", 11, GRAY)
add_bullet(tf, "No rotation", 11, GRAY)
add_bullet(tf, "Reference always present", 11, GRAY)

card2 = add_rect(s, 5.3, 2.85, 4.2, 2.0, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 5.5, 2.9, 3.8, 1.9, "", font_size=12)
add_para(tf, "Phase 2 Extensions", 12, ACCENT_LT, True)
add_bullet(tf, "Zoom unknown — uniform [8×, 12×]", 11, GRAY)
add_bullet(tf, "Rotation unknown — ±5°", 11, GRAY)
add_bullet(tf, "~20% absent pairs (rejection required)", 11, GRAY)
add_bullet(tf, "Degraded images (charging, defocus, noise)", 11, GRAY)
add_bullet(tf, "20s hard timeout, ≤5s median target", 11, GRAY)
add_bullet(tf, "Output: x, y, theta, scale, found, score", 11, GRAY)

# Scoring card
card3 = add_rect(s, 0.5, 5.0, 9.0, 1.8, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 0.7, 5.05, 8.6, 1.7, "", font_size=12)
add_para(tf, "Phase 2 Scoring Rubric (100 pts)", 12, ACCENT_LT, True)
add_para(tf, "40 Localization (tiered, 0.45×A + 0.55×B)  ·  20 Pose (10 scale + 10 rotation)  ·  "
         "15 Rejection F1  ·  10 Calibration AUC  ·  5 Efficiency  ·  10 Carried Forward",
         11, GRAY)

footer_bar(s)

# ═══ SLIDE 3: DATASET GENERATOR ═══════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "02", "SYNTHETIC DATASET", "Physically-Motivated SEM Image Synthesis")

# Left — structure rendering
card = add_rect(s, 0.5, 1.35, 4.3, 2.0, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 0.7, 1.4, 4.0, 1.9, "", font_size=12)
add_para(tf, "Structure Rendering", 13, ACCENT_LT, True)
add_bullet(tf, "DRAM — word-line / bit-line / contact array", 11, GRAY)
add_bullet(tf, "FinFET — parallel fin / gate-line array", 11, GRAY)
add_bullet(tf, "RGB Optical — 3-channel diffraction + chromatic aberration", 11, GRAY)
add_para(tf, "Same renderer at both scales — only pitch/width parameters change.", 10, MUTED)

# Left — LER
card = add_rect(s, 0.5, 3.55, 4.3, 1.6, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 0.7, 3.6, 4.0, 1.5, "", font_size=12)
add_para(tf, "Line-Edge Roughness (LER)", 13, ACCENT_LT, True)
add_para(tf, "Smoothed random perturbation of each edge — spatially correlated, "
         "matching real lithographic processes. Makes neighbouring cells visually similar "
         "but not pixel-identical — source of genuine periodic ambiguity.", 11, GRAY)

# Right — SEM noise model
card = add_rect(s, 5.2, 1.35, 4.3, 2.8, BG_CARD, ACCENT)
_, tf = add_textbox(s, 5.4, 1.4, 4.0, 2.7, "", font_size=12)
add_para(tf, "SEM Noise Model (4 Components)", 13, ACCENT_LT, True)
add_bullet(tf, "Poisson (shot) — signal-dependent, dominant in SEM", 11, GRAY)
add_bullet(tf, "Gaussian (read) — additive, signal-independent", 11, GRAY)
add_bullet(tf, "Beam-spot PSF blur — Gaussian convolution", 11, GRAY)
add_bullet(tf, "Scan-line jitter — per-row intensity offset", 11, GRAY)
add_para(tf, "", 6, GRAY)
add_para(tf, "Search image deliberately noisier (lower gain, larger blur) — "
         "modeling fast survey scan vs. careful characterization.", 10, MUTED)

# Right — Phase 2 degradation
card = add_rect(s, 5.2, 4.35, 4.3, 1.6, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 5.4, 4.4, 4.0, 1.5, "", font_size=12)
add_para(tf, "Phase 2 Degradation Effects", 13, YELLOW, True)
add_bullet(tf, "Charging artifacts (localized bright blooming)", 11, GRAY)
add_bullet(tf, "Scan distortion (smooth 2D displacement field)", 11, GRAY)
add_bullet(tf, "Polygon CD jitter (±20% critical dimension)", 11, GRAY)
add_bullet(tf, "Elevated shot noise + defocus", 11, GRAY)

footer_bar(s)

# ═══ SLIDE 4: LITERATURE GROUNDING ════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "03", "LITERATURE GROUNDING", "Every Design Choice Justified Against Public Sources")

data = [
    ["Design Choice", "Source(s)", "Code Location"],
    ["Poisson-Gaussian SEM noise",
     "Avci et al.; ScienceDirect M-Denoiser 2023; Oxford Academic 2025",
     "apply_sem_noise()"],
    ["Line-Edge Roughness (LER/LWR)",
     "Bunday/Bishop/Villarrubia (NIST/SEMATECH); ITRS targets",
     "_smooth_noise_1d()"],
    ["FinFET pitch/width parameters",
     'IEEE "Scaling of SOI FinFETs" — pitch 40nm, gate 20nm',
     "render_finfet_pattern()"],
    ["Optical chromatic aberration",
     "Nikon MicroscopyU; Wadsworth Center glossary",
     "to_rgb_optical()"],
    ["Fast NCC via FFT",
     'Lewis, "Fast Normalized Cross-Correlation," 1995',
     "cv2.matchTemplate"],
    ["Peak-to-Sidelobe Ratio",
     "Bolme et al. MOSSE tracker, CVPR 2010",
     "peak_to_sidelobe_ratio()"],
    ["Prior-window restriction",
     "US Patent 7,545,497; arXiv 2012.12784; Patent 12,327,739",
     "expected_xy / max_drift_px"],
]
add_table(s, 0.5, 1.35, 9.0, 4.5, 8, 3, data,
          col_widths=[2.8, 4.0, 2.2], font_size=10)

_, tf = add_textbox(s, 0.5, 6.3, 9, 0.4, "", font_size=9)
add_para(tf, "Full annotated bibliography with URLs: docs/citations.md — 15+ verifiable public sources.", 9, MUTED)

footer_bar(s)

# ═══ SLIDE 5: PHASE 1 ALGORITHM ═══════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "04", "PHASE 1 — ALGORITHM", "Classical FFT-NCC Matched Filtering + Prior Window")

# Pipeline boxes
steps = [
    ("INPUT", "Reference\n1000×1000\n1 nm/px"),
    ("STEP 1", "Downsample\nby known 10×\n→ 100×100"),
    ("STEP 2", "FFT-NCC\nMatched Filter\nO(N log N)"),
    ("STEP 3", "Top-K Peaks\nNon-Max\nSuppression"),
    ("INNOVATION", "Prior Window\nFilter by\nexpected coord"),
    ("STEP 4", "Sub-pixel\nParabolic\nInterpolation"),
]
for i, (label, desc) in enumerate(steps):
    x = 0.4 + i * 1.55
    border = GREEN if label == "INNOVATION" else RGBColor(0x2D, 0x3A, 0x50)
    add_rect(s, x, 1.4, 1.35, 1.5, BG_CARD, border)
    lbl_color = GREEN if label == "INNOVATION" else ACCENT
    add_textbox(s, x + 0.05, 1.45, 1.25, 0.25, label,
                font_size=7, color=lbl_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.05, 1.7, 1.25, 1.1, desc,
                font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_textbox(s, x + 1.35, 1.9, 0.2, 0.4, "→",
                    font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Prior Window card
card = add_rect(s, 0.5, 3.2, 4.3, 2.3, BG_CARD, GREEN)
_, tf = add_textbox(s, 0.7, 3.25, 4.0, 2.2, "", font_size=12)
add_para(tf, "Prior-Window Innovation", 13, GREEN, True)
add_para(tf, "The tool knows the commanded stage position. Drift is bounded "
         "(typically <60 px). Discarding candidates outside this window is:", 11, GRAY)
add_bullet(tf, "Faster — smaller effective search area", 11, GRAY)
add_bullet(tf, "More robust — distant periodic decoys eliminated", 11, GRAY)
add_bullet(tf, "Zero precision loss — same FFT-NCC inside window", 11, GRAY)
add_para(tf, "Grounded: US Patent 7,545,497", 9, MUTED)

# GAR card
card = add_rect(s, 5.2, 3.2, 4.3, 2.3, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 5.4, 3.25, 4.0, 2.2, "", font_size=12)
add_para(tf, "Global Ambiguity Ratio (GAR)", 13, ACCENT_LT, True)
add_para(tf, "Corrected confidence metric replacing local PSR. Compares best NCC "
         "peak vs best non-local competing peak (Lowe's ratio test style).", 11, GRAY)
add_para(tf, "", 4, GRAY)
add_bullet(tf, "GAR ≈ 0.88–0.91 → easy (clear unique match)", 11, GREEN)
add_bullet(tf, "GAR ≈ 0.98–0.999 → ambiguous (decoy present)", 11, RED)
add_para(tf, "PSR only sees local sidelobes — wrong scope for distant decoys.", 9, MUTED)

footer_bar(s)

# ═══ SLIDE 6: PHASE 1 RESULTS ═════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "05", "PHASE 1 — RESULTS", "100% Accuracy with Sub-Pixel Precision")

data = [
    ["Configuration", "Success Rate", "Mean Error", "Median Error", "Time"],
    ["Non-adversarial (easy)", ("100% (30/30)", GREEN, True), "0.03 px", "0.03 px", "~309 ms"],
    ["Adversarial decoy, no prior", ("46.7% (14/30)", RED, True), ("194.9 px", RED), "155.9 px", "~224 ms"],
    ["Adversarial decoy + prior", ("100% (30/30)", GREEN, True), ("0.02 px", GREEN, True), ("0.02 px", GREEN), "~215 ms"],
]
add_table(s, 0.5, 1.35, 9.0, 1.5, 4, 5, data,
          col_widths=[3.0, 1.6, 1.3, 1.3, 1.0], font_size=11)

# DL ablation
_, tf = add_textbox(s, 0.5, 3.1, 5, 0.3, "", font_size=14)
add_para(tf, "DL Ablation — Siamese CNN (69K parameters, trained from scratch)", 13, ACCENT_LT, True)

dl_data = [
    ["Configuration", "Success", "Mean Error", "Time"],
    ["Raw DL argmax", "6.7%", "73.6 px", "4.4 ms"],
    ["DL + NCC refinement", "50.0%", "66.1 px", "7.1 ms"],
    ["DL + refinement + prior", ("86–90%", YELLOW, True), "~5–8 px", "10–49 ms"],
]
add_table(s, 0.5, 3.55, 5.0, 1.5, 4, 4, dl_data,
          col_widths=[2.2, 1.0, 1.0, 0.8], font_size=10)

# Failure case card
card = add_rect(s, 5.8, 3.1, 3.7, 2.8, BG_CARD, RED)
_, tf = add_textbox(s, 6.0, 3.15, 3.4, 2.7, "", font_size=12)
add_para(tf, "Honest Failure Case (Seed 24)", 12, RED, True)
add_para(tf, "Error: 753 px (without prior)", 11, YELLOW, True)
add_para(tf, "", 4, GRAY)
add_para(tf, "True site: 0.8464 vs decoy: 0.8438 — 0.3% gap. Algorithm "
         "chose decoy (closer to center per tie-break rule).", 10, GRAY)
add_para(tf, "", 4, GRAY)
add_para(tf, "Not a bug: inherent limitation when sites are structurally "
         "identical. Resolved by prior window.", 10, GRAY)
add_para(tf, "With prior → 0.02 px ✓", 11, GREEN, True)

_, tf = add_textbox(s, 0.5, 5.3, 9, 0.5, "", font_size=10)
add_para(tf, "Classical FFT-NCC submitted as primary method. DL presented as honest "
         "trade-off study — 4–20× faster but hasn't converged to match classical accuracy.", 10, MUTED)

footer_bar(s)

# ═══ SLIDE 7: PHASE 2 ALGORITHM ═══════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "06", "PHASE 2 — ALGORITHM", "Multi-Stage Beam Search over (Scale, θ, Position)")

# Stage cards
stages = [
    ("Stage 1 — Coarse Grid", ACCENT, [
        "9 scales × 5 θ anchors = 45 evals",
        "Scale: {8.0, 8.5, ..., 12.0}",
        "θ: {-5, -2.5, 0, 2.5, 5}°",
        "cv2.matchTemplate(TM_CCOEFF_NORMED)",
        "Output: Top-3 (scale, θ) by score",
    ]),
    ("Stage 2 — θ Refinement", ACCENT, [
        "Per beam candidate:",
        "Sweep θ ±1.5° in 0.5° steps",
        "Pick best θ per anchor",
        "Output: Top-3 after re-sort",
        "",
    ]),
    ("Stage 3 — Fine Joint", ACCENT, [
        "Scale ±0.3 (step 0.3)",
        "θ ±0.6° (step 0.3°)",
        "Multi-peak: top-4 spatial peaks",
        "min_sep = 15% template size",
        "Prevents aliased periodic lobe",
    ]),
]
for i, (title, color, bullets) in enumerate(stages):
    x = 0.5 + i * 3.1
    add_rect(s, x, 1.35, 2.9, 2.6, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
    _, tf = add_textbox(s, x + 0.15, 1.4, 2.6, 2.5, "", font_size=11)
    add_para(tf, title, 12, ACCENT_LT, True)
    for b in bullets:
        if b:
            add_bullet(tf, b, 10, GRAY)

# Final selection card
card = add_rect(s, 0.5, 4.15, 4.5, 1.6, BG_CARD, GREEN)
_, tf = add_textbox(s, 0.7, 4.2, 4.2, 1.5, "", font_size=12)
add_para(tf, "Final Selection", 13, GREEN, True)
add_para(tf, "confidence = score × (1 − GAR)", 14, ACCENT_LT, True, alignment=PP_ALIGN.CENTER)
add_para(tf, "Winner = highest confidence", 11, GRAY)
add_para(tf, "→ Subpixel parabolic refinement of (x, y)", 11, GRAY)

# Rejection card
card = add_rect(s, 5.2, 4.15, 4.3, 1.6, BG_CARD, YELLOW)
_, tf = add_textbox(s, 5.4, 4.2, 4.0, 1.5, "", font_size=12)
add_para(tf, "Rejection Logic", 13, YELLOW, True)
add_para(tf, "found = (score ≥ 0.43) OR (gar ≤ 0.65)", 13, WHITE, True, alignment=PP_ALIGN.CENTER)
add_para(tf, "Score = primary separator. GAR rescues borderline-score "
         "present pairs with high spatial confidence.", 10, GRAY)
add_para(tf, "Calibrated on real 20-pair organizer data. F1 = 1.000.", 10, MUTED)

# Safety
card = add_rect(s, 0.5, 6.0, 9.0, 0.8, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 0.7, 6.05, 8.6, 0.7, "", font_size=11)
add_para(tf, "Safety:  17.5s wall-clock budget guard  ·  Subprocess isolation per pair  ·  "
         "Timeout/crash → found=0  ·  Median actual: 2.13s/pair", 10, GRAY)

footer_bar(s)

# ═══ SLIDE 8: CRITICAL BUG FIX ════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "07", "CRITICAL DISCOVERY", "The Single Biggest Fix: Core NCC Primitive")

# Before card
card = add_rect(s, 0.5, 1.35, 4.3, 2.8, BG_CARD, RED)
_, tf = add_textbox(s, 0.7, 1.4, 4.0, 2.7, "", font_size=12)
add_para(tf, "✗  Hand-Rolled FFT-NCC (Before)", 13, RED, True)
add_para(tf, "Custom integral images + FFT correlation", 11, GRAY)
add_para(tf, "", 4, GRAY)
add_bullet(tf, "Location Error: 18 – 358 px", 11, RED, True)
add_bullet(tf, "NCC Score: 0.66 – 0.82", 11, GRAY)
add_bullet(tf, "Status: Wrong location at high score", 11, RED)
add_para(tf, "", 4, GRAY)
add_para(tf, "Verified on real organizer data at exact GT pose. "
         "Confident but wrong — hardest class of bug.", 10, MUTED)

# After card
card = add_rect(s, 5.2, 1.35, 4.3, 2.8, BG_CARD, GREEN)
_, tf = add_textbox(s, 5.4, 1.4, 4.0, 2.7, "", font_size=12)
add_para(tf, "✓  cv2.matchTemplate (After)", 13, GREEN, True)
add_para(tf, "TM_CCOEFF_NORMED — OpenCV verified NCC", 11, GRAY)
add_para(tf, "", 4, GRAY)
add_bullet(tf, "Location Error: 0.3 – 0.7 px", 11, GREEN, True)
add_bullet(tf, "NCC Score: 0.85 – 0.87", 11, GRAY)
add_bullet(tf, "Status: Correct on all 20 pairs", 11, GREEN)
add_para(tf, "", 4, GRAY)
add_para(tf, "Same inputs, same (scale, θ). Problem was never "
         "search strategy — it was the correlation primitive.", 10, MUTED)

# Other fixes
card = add_rect(s, 0.5, 4.4, 9.0, 1.4, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 0.7, 4.45, 8.6, 1.3, "", font_size=12)
add_para(tf, "Other Key Fixes", 13, ACCENT_LT, True)
add_bullet(tf, "Scale step 1.0 → 0.5 — true scale at 10.55 was 0.45 from nearest grid, causing wrong-position lock", 10, GRAY)
add_bullet(tf, "Absent generator — same-pitch decoy aliases; switched to mismatched-periodicity reference (matching organizer)", 10, GRAY)
add_bullet(tf, "Threshold calibration — score gap 0.0001 between p012/p018; GAR provides 0.13 real margin", 10, GRAY)

footer_bar(s)

# ═══ SLIDE 9: PHASE 2 RESULTS ═════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "08", "PHASE 2 — RESULTS", "Verified on Real Organizer Sample Data (20 Pairs)")

# Metric boxes
metrics = [
    ("20/20", "All Pairs Correct", GREEN),
    ("0.26 px", "Median Error", ACCENT_LT),
    ("F1 = 1.0", "Rejection Score", GREEN),
    ("2.13s", "Median Time/Pair", ACCENT_LT),
]
for i, (val, label, color) in enumerate(metrics):
    x = 0.5 + i * 2.35
    add_rect(s, x, 1.35, 2.1, 0.95, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
    add_textbox(s, x, 1.4, 2.1, 0.5, val,
                font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, 1.85, 2.1, 0.3, label,
                font_size=8, color=MUTED, bold=True, alignment=PP_ALIGN.CENTER)

# Per-set table
set_data = [
    ["Set", "Description", "Accuracy", "Median Err", "Max Err"],
    ["A (8 pairs)", "Nominal pose", ("8/8 (1.000)", GREEN, True), "0.20 px", "0.71 px"],
    ["B (6 pairs)", "Degraded (charging, noise)", ("6/6 (1.000)", GREEN, True), "0.44 px", "1.18 px"],
    ["C (4 pairs)", "Absent — rejection", ("4/4  F1=1.0", GREEN, True), "—", "—"],
    ["D (2 pairs)", "RGB optical bonus", ("2/2 (1.000)", GREEN, True), "0.28 px", "0.31 px"],
]
add_table(s, 0.5, 2.55, 4.8, 2.0, 5, 5, set_data,
          col_widths=[1.0, 1.5, 1.1, 0.7, 0.7], font_size=10)

# vs baseline
vs_data = [
    ["Metric", "Ours", "Organizer Baseline"],
    ["Set A accuracy", ("1.000", GREEN, True), "1.000"],
    [("Set B accuracy", WHITE, True), ("1.000", GREEN, True), ("0.467", RED, True)],
    ["Set D accuracy", ("1.000", GREEN, True), "1.000"],
    [("Rejection F1", WHITE, True), ("1.000", GREEN, True), ("0.897", YELLOW)],
    ["Median time", ("2.13s", GREEN, True), "—"],
]
add_table(s, 5.6, 2.55, 3.9, 2.3, 6, 3, vs_data,
          col_widths=[1.4, 1.1, 1.4], font_size=10)

_, tf = add_textbox(s, 5.6, 5.0, 3.9, 0.5, "", font_size=9)
add_para(tf, "Set B improvement is the strongest differentiator — "
         "degraded images are where beam search dramatically outperforms naive ZNCC.", 9, MUTED)

# Loc tiers
_, tf = add_textbox(s, 0.5, 4.8, 4.8, 0.3, "", font_size=11)
add_para(tf, "Localization Tiers", 11, ACCENT_LT, True)

tier_data = [
    ["Error", "Credit", "Our Results"],
    ["≤ 1 px", ("1.0 (full)", GREEN), ("19/20 pairs", GREEN, True)],
    ["≤ 2 px", "0.8", "1 pair (p010: 1.18px)"],
    ["≤ 3 px", "0.6", "—"],
    ["> 5 px", ("0.0", RED), "—"],
]
add_table(s, 0.5, 5.2, 4.8, 1.5, 5, 3, tier_data,
          col_widths=[0.9, 1.2, 2.7], font_size=9)

footer_bar(s)

# ═══ SLIDE 10: SCORING ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "09", "SCORING BREAKDOWN", "Self-Assessment Against Official Rubric")

score_data = [
    ["Component", "Max Pts", "Estimated", "Notes"],
    ["Localization (0.45×A + 0.55×B)", "40", ("~39.6", GREEN, True), "All tier 1 except one tier 2"],
    ["Pose — Scale recovery", "10", ("~10", GREEN, True), "Scale error <1% on all pairs"],
    ["Pose — Rotation recovery", "10", ("~10", GREEN, True), "θ error <0.25° on all pairs"],
    ["Rejection F1", "15", ("15", GREEN, True), "F1 = 1.000 on sample set"],
    ["Confidence calibration (AUC)", "10", ("~8–10", YELLOW), "Score separates outcomes cleanly"],
    ["Efficiency", "5", ("~4–5", GREEN), "2.13s — well under 5s target"],
    ["Carried forward (generator+citations)", "10", ("~8–10", GREEN), "15+ sources, failure analysis"],
    [("TOTAL", ACCENT_LT, True), ("100", ACCENT_LT, True), ("~94–100", GREEN, True), ""],
]
add_table(s, 0.5, 1.35, 9.0, 3.2, 9, 4, score_data,
          col_widths=[3.5, 1.0, 1.2, 3.3], font_size=11)

# Threshold calibration card
card = add_rect(s, 0.5, 4.8, 9.0, 1.6, BG_CARD, YELLOW)
_, tf = add_textbox(s, 0.7, 4.85, 8.6, 1.5, "", font_size=12)
add_para(tf, "Threshold Calibration — Hardest Case", 13, YELLOW, True)
add_para(tf, "p012 (present, degraded): score = 0.4292   vs.   p018 (absent): score = 0.4291", 11, WHITE, True)
add_para(tf, "Score gap: 0.0001 — essentially identical. But GAR provides real separation:", 11, GRAY)
add_bullet(tf, "p012: GAR = 0.6227 ≤ 0.65 → found = 1 ✓  (0.13 margin)", 11, GREEN, True)
add_bullet(tf, "p018: GAR = 0.7527 > 0.65 → found = 0 ✓", 11, GREEN, True)

footer_bar(s)

# ═══ SLIDE 11: FAILURE ANALYSIS ═══════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "10", "HONEST ASSESSMENT", "Failure Modes, Limitations & Explainability")

# Failure 1
card = add_rect(s, 0.5, 1.35, 4.3, 1.7, BG_CARD, RED)
_, tf = add_textbox(s, 0.7, 1.4, 4.0, 1.6, "", font_size=11)
add_para(tf, "Failure: Periodic-Array Ambiguity", 12, RED, True)
add_para(tf, "Structurally identical sites produce NCC scores within 0.3%. "
         "Image-only tie-break fails when sites are physically indistinguishable.", 10, GRAY)
add_para(tf, "Resolution: Prior window uses stage position → 100% ✓", 10, GREEN, True)

# Failure 2
card = add_rect(s, 5.2, 1.35, 4.3, 1.7, BG_CARD, YELLOW)
_, tf = add_textbox(s, 5.4, 1.4, 4.0, 1.6, "", font_size=11)
add_para(tf, "Failure: Scale Quantization", 12, YELLOW, True)
add_para(tf, "At step=1.0, true scale 10.55 falls 0.45 from grid. Degraded NCC lets "
         "unrelated matches win — confident wrong positions at 81–226 px.", 10, GRAY)
add_para(tf, "Resolution: step=0.5 → all sub-1px ✓", 10, GREEN, True)

# Failure 3
card = add_rect(s, 0.5, 3.3, 4.3, 1.7, BG_CARD, YELLOW)
_, tf = add_textbox(s, 0.7, 3.35, 4.0, 1.6, "", font_size=11)
add_para(tf, "Failure: FFT-NCC Implementation Bug", 12, YELLOW, True)
add_para(tf, "Hand-rolled NCC returned wrong locations (18–358 px) at high confidence. "
         "Algorithm appeared correct but silently produced wrong results.", 10, GRAY)
add_para(tf, "Resolution: cv2.matchTemplate → 0.3–0.7 px ✓", 10, GREEN, True)

# Known limitations
card = add_rect(s, 5.2, 3.3, 4.3, 1.7, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 5.4, 3.35, 4.0, 1.6, "", font_size=11)
add_para(tf, "Known Remaining Limitations", 12, ACCENT_LT, True)
add_bullet(tf, "Borderline score gap (0.0001) on hardest pair", 10, GRAY)
add_bullet(tf, "RGB handled as grayscale — no color matching", 10, GRAY)
add_bullet(tf, "Local calibration doesn't fully match organizer data", 10, GRAY)
add_bullet(tf, "Efficiency rank depends on other teams", 10, GRAY)

# Takeaway
card = add_rect(s, 0.5, 5.3, 9.0, 1.0, BG_CARD, GREEN)
_, tf = add_textbox(s, 0.7, 5.35, 8.6, 0.9, "", font_size=12)
add_para(tf, "Key Takeaway", 13, GREEN, True)
add_para(tf, "Every failure mode discovered was documented honestly and resolved systematically. "
         "The prior-window concept, the NCC primitive switch, and threshold calibration via GAR "
         "each emerged directly from analyzing specific, reproducible failure cases.", 11, GRAY)

footer_bar(s)

# ═══ SLIDE 12: DELIVERABLES ═══════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "11", "DELIVERABLES", "Code, Documentation & Reproducibility")

# Phase 2 code
card = add_rect(s, 0.5, 1.35, 4.3, 2.2, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 0.7, 1.4, 4.0, 2.1, "", font_size=11)
add_para(tf, "Phase 2 Codebase", 12, ACCENT_LT, True)
add_bullet(tf, "register.py — CLI entry point", 10, GRAY)
add_bullet(tf, "localize_v2.py — Beam-search NCC localizer", 10, GRAY)
add_bullet(tf, "generate_dataset_v2.py — Extended generator", 10, GRAY)
add_bullet(tf, "scorer.py — Self-assessment scorer", 10, GRAY)
add_bullet(tf, "calibrate_*.py / sweep_thresholds.py", 10, GRAY)

# Phase 1 code
card = add_rect(s, 5.2, 1.35, 4.3, 2.2, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 5.4, 1.4, 4.0, 2.1, "", font_size=11)
add_para(tf, "Phase 1 Codebase (Preserved)", 12, ACCENT_LT, True)
add_bullet(tf, "generate_dataset.py — DRAM/FinFET/RGB generator", 10, GRAY)
add_bullet(tf, "localize.py — Classical FFT-NCC localizer", 10, GRAY)
add_bullet(tf, "evaluate.py — Batch evaluation harness", 10, GRAY)
add_bullet(tf, "dl_localize.py — Siamese CNN ablation (69K params)", 10, GRAY)

# Docs
card = add_rect(s, 0.5, 3.8, 4.3, 1.4, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 0.7, 3.85, 4.0, 1.3, "", font_size=11)
add_para(tf, "Documentation", 12, ACCENT_LT, True)
add_bullet(tf, "docs/proposal.pdf — Full technical proposal", 10, GRAY)
add_bullet(tf, "docs/citations.md — 15+ annotated sources", 10, GRAY)
add_bullet(tf, "README.md — Setup & run instructions", 10, GRAY)

# Environment
card = add_rect(s, 5.2, 3.8, 4.3, 1.4, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
_, tf = add_textbox(s, 5.4, 3.85, 4.0, 1.3, "", font_size=11)
add_para(tf, "Environment", 12, ACCENT_LT, True)
add_bullet(tf, "Python 3.12, NumPy 2.5, SciPy 1.18, OpenCV 5.0", 10, GRAY)
add_bullet(tf, "No GPU, no deep learning at inference", 10, GRAY)
add_bullet(tf, "No network access required", 10, GRAY)
add_bullet(tf, "4-core CPU, 8 GB RAM", 10, GRAY)

# GitHub
card = add_rect(s, 0.5, 5.5, 9.0, 1.0, BG_CARD, ACCENT)
_, tf = add_textbox(s, 0.7, 5.55, 8.6, 0.9, "", font_size=14)
add_para(tf, "GitHub Repository", 14, ACCENT_LT, True)
add_para(tf, "github.com/YashKothari06/Drift-sense", 16, WHITE, True)
add_para(tf, "All Phase 1 + Phase 2 code, documentation, sample outputs, and trained weights.", 10, MUTED)

footer_bar(s)

# ═══ SLIDE 13: CONCLUSION ═════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_DARK)

# Badge
badge = add_rect(s, 3.5, 1.2, 3.0, 0.35, GREEN)
add_textbox(s, 3.5, 1.22, 3.0, 0.3, "PHASE 2 COMPLETE",
            font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Title
add_textbox(s, 1, 1.9, 8, 1.0,
            "Sub-Pixel Accuracy.\nReal-Time Speed.\nIndustry-Grounded.",
            font_size=30, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Metrics
metrics = [
    ("100%", "Sample\nAccuracy", GREEN),
    ("0.26 px", "Median\nError", ACCENT_LT),
    ("F1 = 1.0", "Rejection\nScore", GREEN),
    ("2.13s", "Per Pair\nSpeed", ACCENT_LT),
    ("15+", "Literature\nCitations", ACCENT_LT),
]
for i, (val, label, color) in enumerate(metrics):
    x = 0.6 + i * 1.85
    add_rect(s, x, 3.3, 1.6, 1.2, BG_CARD, RGBColor(0x2D, 0x3A, 0x50))
    add_textbox(s, x, 3.35, 1.6, 0.6, val,
                font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, 3.95, 1.6, 0.5, label,
                font_size=8, color=MUTED, bold=True, alignment=PP_ALIGN.CENTER)

# Divider
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.4), Inches(4.8), Inches(1.2), Pt(2))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
line.shadow.inherit = False

# Summary line
add_textbox(s, 1, 5.1, 8, 0.8,
            "Classical FFT-NCC + multi-stage beam search + cv2.matchTemplate\n"
            "Beating the organizer baseline on all scoring axes.",
            font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# Author
add_textbox(s, 1, 6.0, 8, 0.4, "Yash Kothari  ·  BITS Pilani  ·  github.com/YashKothari06",
            font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1, 6.4, 8, 0.4, "Thank You",
            font_size=16, color=ACCENT_LT, bold=True, alignment=PP_ALIGN.CENTER)

footer_bar(s)

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "drift_sense_presentation.pptx")
prs.save(out_path)
print(f"✓ Saved presentation to: {out_path}")
print(f"  {len(prs.slides)} slides generated.")
print(f"\nTo convert to PDF:")
print(f"  • Open in PowerPoint/LibreOffice → File → Export as PDF")
print(f"  • Or: libreoffice --headless --convert-to pdf drift_sense_presentation.pptx")
